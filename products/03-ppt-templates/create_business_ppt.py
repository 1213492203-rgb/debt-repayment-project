"""
Generate a professional 10-slide business plan PowerPoint template.
Blue/white color scheme, Chinese language, editable text.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Color palette
PRIMARY_BLUE = RGBColor(0x1A, 0x56, 0xDB)
DARK_BLUE = RGBColor(0x0F, 0x2B, 0x5B)
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xFF)
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1F, 0x29, 0x3A)
GRAY_TEXT = RGBColor(0x6B, 0x72, 0x80)
ACCENT_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]  # Blank


def add_bg_shape(slide, color=WHITE):
    """Add a full-slide background rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_top_bar(slide):
    """Add a thin decorative bar at the top."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, Inches(0.06)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY_BLUE
    bar.line.fill.background()


def add_bottom_bar(slide):
    """Add footer bar."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.05),
        prs.slide_width, Inches(0.45)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()
    # Footer text
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "商业计划书  |  机密文件"
    p.font.size = Pt(9)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.RIGHT
    tf.margin_right = Inches(0.5)
    tf.margin_top = Pt(4)


def add_slide_number(slide, num):
    """Add slide number indicator."""
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.95), Inches(1), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num:02d}"
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.font.bold = True


def add_section_title(slide, title, subtitle=None):
    """Add a styled title to a content slide."""
    # Left accent bar
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.35),
        Inches(0.06), Inches(0.6)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = PRIMARY_BLUE
    accent.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.3), Inches(10), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.color.rgb = DARK_BLUE
    p.font.bold = True

    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(0.75), Inches(0.85), Inches(10), Inches(0.4))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY_TEXT


def add_card_shape(slide, left, top, width, height, fill_color=WHITE, border_color=None):
    """Add a rounded rectangle card."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


# ==========================================
# SLIDE 1: COVER
# ==========================================
slide1 = prs.slides.add_slide(blank_layout)

# Full blue background
bg = slide1.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
    prs.slide_width, prs.slide_height
)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK_BLUE
bg.line.fill.background()

# Decorative geometric shapes
circle1 = slide1.shapes.add_shape(
    MSO_SHAPE.OVAL, Inches(9.5), Inches(-1.5),
    Inches(5.5), Inches(5.5)
)
circle1.fill.solid()
circle1.fill.fore_color.rgb = RGBColor(0x15, 0x35, 0x6E)
circle1.line.fill.background()
circle1.rotation = 15

circle2 = slide1.shapes.add_shape(
    MSO_SHAPE.OVAL, Inches(-1), Inches(4.5),
    Inches(4), Inches(4)
)
circle2.fill.solid()
circle2.fill.fore_color.rgb = RGBColor(0x12, 0x2D, 0x58)
circle2.line.fill.background()

# Main title
txBox = slide1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(9), Inches(1.2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "商业计划书"
p.font.size = Pt(56)
p.font.color.rgb = WHITE
p.font.bold = True

# Subtitle
txBox2 = slide1.shapes.add_textbox(Inches(1.2), Inches(2.95), Inches(9), Inches(0.8))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "BUSINESS PLAN"
p2.font.size = Pt(22)
p2.font.color.rgb = ACCENT_BLUE
p2.font.italic = True

# Company name placeholder
txBox3 = slide1.shapes.add_textbox(Inches(1.2), Inches(4.0), Inches(8), Inches(0.6))
tf3 = txBox3.text_frame
p3 = tf3.paragraphs[0]
p3.text = "公司名称  |  日期"
p3.font.size = Pt(18)
p3.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)

# Bottom accent line
line = slide1.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(3.8),
    Inches(2.5), Inches(0.04)
)
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_ORANGE
line.line.fill.background()

# Confidential badge
badge = slide1.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(5.2),
    Inches(2), Inches(0.5)
)
badge.fill.solid()
badge.fill.fore_color.rgb = RGBColor(0x15, 0x35, 0x6E)
badge.line.color.rgb = RGBColor(0x3B, 0x6B, 0xC4)
badge.line.width = Pt(1)
tf_badge = badge.text_frame
tf_badge.word_wrap = True
p_badge = tf_badge.paragraphs[0]
p_badge.text = "🔒  机密文件"
p_badge.font.size = Pt(11)
p_badge.font.color.rgb = ACCENT_BLUE
p_badge.alignment = PP_ALIGN.CENTER


# ==========================================
# SLIDE 2: AGENDA
# ==========================================
slide2 = prs.slides.add_slide(blank_layout)
add_bg_shape(slide2, WHITE)
add_top_bar(slide2)
add_bottom_bar(slide2)
add_slide_number(slide2, 2)
add_section_title(slide2, "目  录", "CONTENTS")

# Agenda items
agenda_items = [
    ("01", "项目概述", "项目背景、愿景与核心价值"),
    ("02", "市场分析", "行业趋势、目标市场与竞争格局"),
    ("03", "产品与解决方案", "核心产品、技术架构与竞争优势"),
    ("04", "商业模式", "收入来源、定价策略与盈利预测"),
    ("05", "团队介绍", "核心团队背景与组织架构"),
    ("06", "财务规划", "融资需求、财务预测与投资回报"),
    ("07", "发展规划", "里程碑、扩张策略与风险应对"),
]

for i, (num, title, desc) in enumerate(agenda_items):
    y = Inches(1.5) + Inches(i * 0.72)

    # Number circle
    circle = slide2.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(1.0), y, Inches(0.45), Inches(0.45)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = PRIMARY_BLUE if i < 3 else ACCENT_BLUE
    circle.line.fill.background()
    tf_c = circle.text_frame
    p_c = tf_c.paragraphs[0]
    p_c.text = num
    p_c.font.size = Pt(14)
    p_c.font.color.rgb = WHITE
    p_c.font.bold = True
    p_c.alignment = PP_ALIGN.CENTER

    # Title
    txBox = slide2.shapes.add_textbox(Inches(1.7), y, Inches(4), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_TEXT
    p.font.bold = True

    # Description
    txBox2 = slide2.shapes.add_textbox(Inches(1.7), y + Inches(0.28), Inches(6), Inches(0.25))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = GRAY_TEXT

# Right side decorative block
deco = slide2.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(10.5), Inches(1.2),
    Inches(2.33), Inches(5.5)
)
deco.fill.solid()
deco.fill.fore_color.rgb = LIGHT_BLUE
deco.line.fill.background()


# ==========================================
# SLIDE 3: DATA CHART
# ==========================================
slide3 = prs.slides.add_slide(blank_layout)
add_bg_shape(slide3, WHITE)
add_top_bar(slide3)
add_bottom_bar(slide3)
add_slide_number(slide3, 3)
add_section_title(slide3, "市场数据分析", "Market Data & Key Metrics")

# KPI Cards
kpi_data = [
    ("¥ 50亿+", "可触达市场规模", "TAM"),
    ("35%", "年复合增长率", "CAGR"),
    ("200万+", "目标用户数", "Target Users"),
    ("¥ 850", "客单价 (ARPU)", "ARPU"),
]

for i, (value, label, tag) in enumerate(kpi_data):
    x = Inches(0.8 + i * 3.1)
    card = add_card_shape(slide3, x, Inches(1.6), Inches(2.8), Inches(2.0),
                          fill_color=WHITE, border_color=RGBColor(0xE5, 0xE7, 0xEB))

    # Tag
    tag_box = slide3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.15), Inches(1.75),
        Inches(0.8), Inches(0.3)
    )
    tag_box.fill.solid()
    tag_box.fill.fore_color.rgb = LIGHT_BLUE
    tag_box.line.fill.background()
    tf_tag = tag_box.text_frame
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = tag
    p_tag.font.size = Pt(8)
    p_tag.font.color.rgb = PRIMARY_BLUE
    p_tag.alignment = PP_ALIGN.CENTER

    # Value
    txBox = slide3.shapes.add_textbox(x + Inches(0.2), Inches(2.2), Inches(2.4), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(28)
    p.font.color.rgb = DARK_BLUE
    p.font.bold = True

    # Label
    txBox2 = slide3.shapes.add_textbox(x + Inches(0.2), Inches(2.8), Inches(2.4), Inches(0.35))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = label
    p2.font.size = Pt(11)
    p2.font.color.rgb = GRAY_TEXT

# Chart placeholder area
chart_area = add_card_shape(slide3, Inches(0.8), Inches(4.0), Inches(11.7), Inches(2.7),
                            fill_color=LIGHT_BLUE, border_color=RGBColor(0xC7, 0xD2, 0xFE))
txBox_c = slide3.shapes.add_textbox(Inches(2), Inches(5.0), Inches(9), Inches(0.5))
tf_c = txBox_c.text_frame
p_c = tf_c.paragraphs[0]
p_c.text = "📊  在此处插入数据图表（柱状图 / 折线图 / 饼图）"
p_c.font.size = Pt(16)
p_c.font.color.rgb = PRIMARY_BLUE
p_c.alignment = PP_ALIGN.CENTER


# ==========================================
# SLIDE 4: TIMELINE
# ==========================================
slide4 = prs.slides.add_slide(blank_layout)
add_bg_shape(slide4, WHITE)
add_top_bar(slide4)
add_bottom_bar(slide4)
add_slide_number(slide4, 4)
add_section_title(slide4, "项目里程碑", "Project Timeline & Milestones")

timeline_events = [
    ("2026 Q1", "项目启动", "完成市场调研与\n产品原型设计"),
    ("2026 Q2", "MVP 上线", "核心功能开发完成\n首批种子用户获取"),
    ("2026 Q3", "产品迭代", "用户反馈驱动优化\n付费转化率突破 5%"),
    ("2026 Q4", "规模增长", "用户突破 10 万\n启动 A 轮融资"),
    ("2027 Q1", "市场扩张", "拓展 3 个新城市\n团队扩充至 30 人"),
]

# Horizontal line
line = slide4.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(3.5),
    Inches(11.3), Inches(0.04)
)
line.fill.solid()
line.fill.fore_color.rgb = PRIMARY_BLUE
line.line.fill.background()

for i, (date, title, desc) in enumerate(timeline_events):
    x = Inches(1.1 + i * 2.3)

    # Node circle
    dot = slide4.shapes.add_shape(
        MSO_SHAPE.OVAL, x + Inches(0.95), Inches(3.3),
        Inches(0.35), Inches(0.35)
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = PRIMARY_BLUE if i < 3 else ACCENT_ORANGE
    dot.line.fill.background()

    card_y = Inches(1.5) if i % 2 == 0 else Inches(4.0)

    # Card
    card = add_card_shape(slide4, x, card_y, Inches(2.1), Inches(1.8),
                          fill_color=WHITE, border_color=RGBColor(0xE5, 0xE7, 0xEB))

    # Date
    txBox = slide4.shapes.add_textbox(x + Inches(0.15), card_y + Inches(0.15),
                                       Inches(1.8), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = date
    p.font.size = Pt(10)
    p.font.color.rgb = PRIMARY_BLUE
    p.font.bold = True

    # Title
    txBox2 = slide4.shapes.add_textbox(x + Inches(0.15), card_y + Inches(0.45),
                                        Inches(1.8), Inches(0.35))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = title
    p2.font.size = Pt(14)
    p2.font.color.rgb = DARK_TEXT
    p2.font.bold = True

    # Description
    txBox3 = slide4.shapes.add_textbox(x + Inches(0.15), card_y + Inches(0.85),
                                        Inches(1.8), Inches(0.8))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = desc
    p3.font.size = Pt(10)
    p3.font.color.rgb = GRAY_TEXT
    p3.line_spacing = Pt(16)


# ==========================================
# SLIDE 5: TEAM
# ==========================================
slide5 = prs.slides.add_slide(blank_layout)
add_bg_shape(slide5, WHITE)
add_top_bar(slide5)
add_bottom_bar(slide5)
add_slide_number(slide5, 5)
add_section_title(slide5, "核心团队", "Meet Our Team")

team_members = [
    ("张伟", "CEO / 创始人", "10年互联网行业经验\n前阿里巴巴产品总监\n北京大学 MBA"),
    ("李娜", "CTO", "15年技术研发经验\n前腾讯高级架构师\n清华计算机博士"),
    ("王磊", "COO", "8年运营管理经验\n前美团区域负责人\n复旦 EMBA"),
    ("陈静", "CMO", "12年市场营销经验\n前字节跳动品牌总监\n海外留学背景"),
]

for i, (name, role, bio) in enumerate(team_members):
    x = Inches(0.7 + i * 3.15)

    # Avatar placeholder
    avatar = slide5.shapes.add_shape(
        MSO_SHAPE.OVAL, x + Inches(0.7), Inches(1.6),
        Inches(1.4), Inches(1.4)
    )
    avatar.fill.solid()
    avatar.fill.fore_color.rgb = LIGHT_BLUE
    avatar.line.color.rgb = PRIMARY_BLUE
    avatar.line.width = Pt(2)

    # Initials
    tf_av = avatar.text_frame
    p_av = tf_av.paragraphs[0]
    p_av.text = name[0]
    p_av.font.size = Pt(32)
    p_av.font.color.rgb = PRIMARY_BLUE
    p_av.font.bold = True
    p_av.alignment = PP_ALIGN.CENTER

    # Name
    txBox = slide5.shapes.add_textbox(x, Inches(3.2), Inches(2.8), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_TEXT
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Role
    txBox2 = slide5.shapes.add_textbox(x, Inches(3.55), Inches(2.8), Inches(0.3))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = role
    p2.font.size = Pt(11)
    p2.font.color.rgb = PRIMARY_BLUE
    p2.alignment = PP_ALIGN.CENTER

    # Bio
    txBox3 = slide5.shapes.add_textbox(x, Inches(3.95), Inches(2.8), Inches(1.5))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = bio
    p3.font.size = Pt(10)
    p3.font.color.rgb = GRAY_TEXT
    p3.alignment = PP_ALIGN.CENTER
    p3.line_spacing = Pt(15)


# ==========================================
# SLIDE 6: SWOT ANALYSIS
# ==========================================
slide6 = prs.slides.add_slide(blank_layout)
add_bg_shape(slide6, WHITE)
add_top_bar(slide6)
add_bottom_bar(slide6)
add_slide_number(slide6, 6)
add_section_title(slide6, "SWOT 分析", "Strategic Analysis")

swot_data = [
    ("S 优势", "Strengths", RGBColor(0x10, 0xB9, 0x81), [
        "核心技术壁垒高",
        "团队执行能力强",
        "先发市场优势",
    ]),
    ("W 劣势", "Weaknesses", RGBColor(0xEF, 0x44, 0x44), [
        "品牌知名度较低",
        "初期资金有限",
        "渠道建设需时间",
    ]),
    ("O 机会", "Opportunities", RGBColor(0xF5, 0x9E, 0x0B), [
        "市场高速增长期",
        "政策红利支持",
        "技术升级窗口期",
    ]),
    ("T 威胁", "Threats", RGBColor(0x8B, 0x5C, 0xF6), [
        "巨头可能入场",
        "技术快速迭代",
        "人才竞争激烈",
    ]),
]

for i, (title, sub, color, items) in enumerate(swot_data):
    col = i % 2
    row = i // 2
    x = Inches(0.7 + col * 6.2)
    y = Inches(1.5 + row * 2.85)

    card = add_card_shape(slide6, x, y, Inches(5.9), Inches(2.6),
                          fill_color=WHITE, border_color=RGBColor(0xE5, 0xE7, 0xEB))

    # Header bar
    header = slide6.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, y, Inches(5.9), Inches(0.55)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = color
    header.line.fill.background()

    tf_h = header.text_frame
    p_h = tf_h.paragraphs[0]
    p_h.text = f"{title}  —  {sub}"
    p_h.font.size = Pt(14)
    p_h.font.color.rgb = WHITE
    p_h.font.bold = True
    p_h.alignment = PP_ALIGN.CENTER

    # Items
    for j, item in enumerate(items):
        txBox = slide6.shapes.add_textbox(x + Inches(0.5), y + Inches(0.75 + j * 0.55),
                                           Inches(4.9), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"•  {item}"
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_TEXT


# ==========================================
# SLIDE 7: FINANCIAL SUMMARY
# ==========================================
slide7 = prs.slides.add_slide(blank_layout)
add_bg_shape(slide7, WHITE)
add_top_bar(slide7)
add_bottom_bar(slide7)
add_slide_number(slide7, 7)
add_section_title(slide7, "财务规划", "Financial Projections")

# Finance cards
fin_data = [
    ("融资需求", "¥ 500 万", "天使轮", "用于产品研发与市场推广"),
    ("预计收入 (Year 1)", "¥ 800 万", "保守预估", "基于当前市场增速"),
    ("毛利率", "65%", "目标值", "行业平均 55-60%"),
    ("投资回报期", "18 个月", "预估", "含 6 个月缓冲"),
]

for i, (label, value, tag, note) in enumerate(fin_data):
    y = Inches(1.5 + i * 1.15)

    # Label
    txBox = slide7.shapes.add_textbox(Inches(1.0), y, Inches(2.5), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(13)
    p.font.color.rgb = GRAY_TEXT

    # Value
    txBox2 = slide7.shapes.add_textbox(Inches(3.5), y, Inches(3), Inches(0.45))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = value
    p2.font.size = Pt(24)
    p2.font.color.rgb = DARK_BLUE
    p2.font.bold = True

    # Tag
    tag_box = slide7.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), y + Inches(0.05),
        Inches(1.2), Inches(0.3)
    )
    tag_box.fill.solid()
    tag_box.fill.fore_color.rgb = LIGHT_BLUE
    tag_box.line.fill.background()
    tf_t = tag_box.text_frame
    p_t = tf_t.paragraphs[0]
    p_t.text = tag
    p_t.font.size = Pt(9)
    p_t.font.color.rgb = PRIMARY_BLUE
    p_t.alignment = PP_ALIGN.CENTER

    # Note
    txBox3 = slide7.shapes.add_textbox(Inches(3.5), y + Inches(0.5), Inches(5), Inches(0.3))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = note
    p3.font.size = Pt(10)
    p3.font.color.rgb = GRAY_TEXT

    # Separator line
    if i < 3:
        sep = slide7.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1.0), y + Inches(0.95),
            Inches(11.3), Inches(0.01)
        )
        sep.fill.solid()
        sep.fill.fore_color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
        sep.line.fill.background()

# Chart placeholder
chart_box = add_card_shape(slide7, Inches(8.5), Inches(5.0), Inches(4.3), Inches(1.8),
                           fill_color=LIGHT_BLUE)
txBox_c = slide7.shapes.add_textbox(Inches(8.8), Inches(5.7), Inches(3.7), Inches(0.5))
tf_c = txBox_c.text_frame
p_c = tf_c.paragraphs[0]
p_c.text = "📈  插入收入增长趋势图"
p_c.font.size = Pt(12)
p_c.font.color.rgb = PRIMARY_BLUE
p_c.alignment = PP_ALIGN.CENTER


# ==========================================
# SLIDE 8: ROADMAP
# ==========================================
slide8 = prs.slides.add_slide(blank_layout)
add_bg_shape(slide8, WHITE)
add_top_bar(slide8)
add_bottom_bar(slide8)
add_slide_number(slide8, 8)
add_section_title(slide8, "发展战略路线图", "Strategic Roadmap")

phases = [
    ("Phase 1", "夯实基础", "2026 H1", [
        "核心产品 MVP 上线",
        "获取 1000 种子用户",
        "建立技术壁垒",
    ], PRIMARY_BLUE),
    ("Phase 2", "快速扩张", "2026 H2", [
        "用户增长至 10 万+",
        "推出付费版本",
        "组建销售团队",
    ], ACCENT_BLUE),
    ("Phase 3", "生态构建", "2027", [
        "开放 API 接口",
        "建立合作伙伴网络",
        "探索海外市场",
    ], ACCENT_ORANGE),
]

for i, (phase, title, period, items, color) in enumerate(phases):
    x = Inches(0.8 + i * 4.1)

    # Phase header
    header = slide8.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, Inches(1.5), Inches(3.8), Inches(0.55)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = color
    header.line.fill.background()
    tf_h = header.text_frame
    p_h = tf_h.paragraphs[0]
    p_h.text = f"{phase}  |  {title}"
    p_h.font.size = Pt(15)
    p_h.font.color.rgb = WHITE
    p_h.font.bold = True
    p_h.alignment = PP_ALIGN.CENTER

    # Period
    txBox = slide8.shapes.add_textbox(x + Inches(0.2), Inches(2.2), Inches(3.4), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = period
    p.font.size = Pt(11)
    p.font.color.rgb = color
    p.font.bold = True

    # Items
    for j, item in enumerate(items):
        txBox2 = slide8.shapes.add_textbox(x + Inches(0.2), Inches(2.6 + j * 0.55),
                                            Inches(3.4), Inches(0.4))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f"{j+1}.  {item}"
        p2.font.size = Pt(12)
        p2.font.color.rgb = DARK_TEXT

# Arrow connectors between phases
for i in range(2):
    x = Inches(4.7 + i * 4.1)
    arrow = slide8.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, x, Inches(1.6), Inches(0.4), Inches(0.3)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    arrow.line.fill.background()


# ==========================================
# SLIDE 9: COMPETITIVE ANALYSIS
# ==========================================
slide9 = prs.slides.add_slide(blank_layout)
add_bg_shape(slide9, WHITE)
add_top_bar(slide9)
add_bottom_bar(slide9)
add_slide_number(slide9, 9)
add_section_title(slide9, "竞争格局分析", "Competitive Landscape")

# Table-like layout
headers = ["维度", "我们", "竞品 A", "竞品 B"]
rows_data = [
    ["产品成熟度", "⭐⭐⭐⭐", "⭐⭐⭐", "⭐⭐⭐⭐⭐"],
    ["价格竞争力", "⭐⭐⭐⭐⭐", "⭐⭐⭐", "⭐⭐"],
    ["技术壁垒", "⭐⭐⭐⭐", "⭐⭐", "⭐⭐⭐⭐"],
    ["用户体验", "⭐⭐⭐⭐⭐", "⭐⭐⭐", "⭐⭐⭐"],
    ["市场占有率", "⭐⭐", "⭐⭐⭐⭐⭐", "⭐⭐⭐⭐"],
]

# Table header
for j, header_text in enumerate(headers):
    x = Inches(1.0 + j * 3.0)
    cell = slide9.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, Inches(1.5), Inches(2.8), Inches(0.55)
    )
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK_BLUE if j == 0 else PRIMARY_BLUE
    cell.line.fill.background()
    tf_c = cell.text_frame
    p_c = tf_c.paragraphs[0]
    p_c.text = header_text
    p_c.font.size = Pt(14)
    p_c.font.color.rgb = WHITE
    p_c.font.bold = True
    p_c.alignment = PP_ALIGN.CENTER

# Table rows
for i, row in enumerate(rows_data):
    for j, cell_text in enumerate(row):
        x = Inches(1.0 + j * 3.0)
        y = Inches(2.05 + i * 0.65)
        cell = slide9.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, Inches(2.8), Inches(0.65)
        )
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC) if i % 2 == 0 else WHITE
        cell.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
        cell.line.width = Pt(0.5)
        tf_c = cell.text_frame
        p_c = tf_c.paragraphs[0]
        p_c.text = cell_text
        p_c.font.size = Pt(12)
        p_c.font.color.rgb = DARK_TEXT if j > 0 else RGBColor(0x1A, 0x56, 0xDB)
        p_c.font.bold = (j == 1)
        p_c.alignment = PP_ALIGN.CENTER

# Insight box
insight_box = add_card_shape(slide9, Inches(1.0), Inches(5.5), Inches(11.3), Inches(1.2),
                             fill_color=RGBColor(0xEF, 0xF6, 0xFF))
txBox_i = slide9.shapes.add_textbox(Inches(1.5), Inches(5.65), Inches(10.3), Inches(0.9))
tf_i = txBox_i.text_frame
tf_i.word_wrap = True
p_i = tf_i.paragraphs[0]
p_i.text = "💡 核心洞察：我们的最大优势在于价格竞争力与用户体验，但需要快速提升市场占有率。建议通过差异化定位和口碑传播实现突围。"
p_i.font.size = Pt(12)
p_i.font.color.rgb = DARK_BLUE
p_i.line_spacing = Pt(20)


# ==========================================
# SLIDE 10: THANK YOU
# ==========================================
slide10 = prs.slides.add_slide(blank_layout)

# Full dark blue background
bg10 = slide10.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
    prs.slide_width, prs.slide_height
)
bg10.fill.solid()
bg10.fill.fore_color.rgb = DARK_BLUE
bg10.line.fill.background()

# Decorative elements
circle10 = slide10.shapes.add_shape(
    MSO_SHAPE.OVAL, Inches(9.5), Inches(-1.5),
    Inches(5.5), Inches(5.5)
)
circle10.fill.solid()
circle10.fill.fore_color.rgb = RGBColor(0x15, 0x35, 0x6E)
circle10.line.fill.background()

# Main text
txBox = slide10.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10), Inches(1.2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "感谢聆听"
p.font.size = Pt(60)
p.font.color.rgb = WHITE
p.font.bold = True

# Sub text
txBox2 = slide10.shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(10), Inches(0.6))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "THANK YOU"
p2.font.size = Pt(24)
p2.font.color.rgb = ACCENT_BLUE
p2.font.italic = True

# Contact info
txBox3 = slide10.shapes.add_textbox(Inches(1.5), Inches(4.3), Inches(8), Inches(1.2))
tf3 = txBox3.text_frame
tf3.word_wrap = True

p3 = tf3.paragraphs[0]
p3.text = "联系方式"
p3.font.size = Pt(18)
p3.font.color.rgb = WHITE
p3.font.bold = True

p4 = tf3.add_paragraph()
p4.text = "邮箱：contact@company.com"
p4.font.size = Pt(14)
p4.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
p4.space_before = Pt(12)

p5 = tf3.add_paragraph()
p5.text = "电话：400-XXX-XXXX"
p5.font.size = Pt(14)
p5.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
p5.space_before = Pt(6)

p6 = tf3.add_paragraph()
p6.text = "地址：北京市朝阳区 XXX 路 XXX 号"
p6.font.size = Pt(14)
p6.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
p6.space_before = Pt(6)

# Accent line
accent10 = slide10.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(4.05),
    Inches(1.8), Inches(0.04)
)
accent10.fill.solid()
accent10.fill.fore_color.rgb = ACCENT_ORANGE
accent10.line.fill.background()


# ==========================================
# SAVE
# ==========================================
output_path = os.path.join(os.path.dirname(__file__), "商业计划书模板.pptx")
prs.save(output_path)
print(f"✅ Successfully generated: {output_path}")
print(f"   Slides: {len(prs.slides)}")
