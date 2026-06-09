"""
Generate a professional 8-slide Work Report PowerPoint template.
Green/white color scheme, Chinese language, editable text.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
GREEN = RGBColor(0x05, 0x9A, 0x69)
DARK_GREEN = RGBColor(0x06, 0x47, 0x3A)
LIGHT_GREEN = RGBColor(0xD1, 0xFA, 0xE5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1F, 0x29, 0x3A)
GRAY = RGBColor(0x6B, 0x72, 0x80)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def add_bg(slide, color=WHITE):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

def add_footer(slide, page_num):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.05), prs.slide_width, Inches(0.45))
    bar.fill.solid(); bar.fill.fore_color.rgb = DARK_GREEN; bar.line.fill.background()
    t = slide.shapes.add_textbox(Inches(0.5), Inches(6.95), Inches(2), Inches(0.4))
    p = t.text_frame.paragraphs[0]; p.text = f"{page_num:02d}"; p.font.size = Pt(10)
    p.font.color.rgb = WHITE; p.font.bold = True

def add_title(slide, title, subtitle=""):
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.3), Inches(0.05), Inches(0.55))
    accent.fill.solid(); accent.fill.fore_color.rgb = GREEN; accent.line.fill.background()
    t = slide.shapes.add_textbox(Inches(0.75), Inches(0.25), Inches(10), Inches(0.55))
    p = t.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(26); p.font.color.rgb = DARK; p.font.bold = True
    if subtitle:
        t2 = slide.shapes.add_textbox(Inches(0.75), Inches(0.75), Inches(10), Inches(0.3))
        p2 = t2.text_frame.paragraphs[0]; p2.text = subtitle; p2.font.size = Pt(12); p2.font.color.rgb = GRAY

def add_card(slide, left, top, w, h, fill=WHITE, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border: s.line.color.rgb = border; s.line.width = Pt(0.5)
    else: s.line.fill.background()
    return s

# === SLIDE 1: Cover ===
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = DARK_GREEN; bg.line.fill.background()
# Deco circles
c1 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-1.5), Inches(5), Inches(5))
c1.fill.solid(); c1.fill.fore_color.rgb = RGBColor(0x05, 0x5A, 0x3E); c1.line.fill.background()
c2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(4), Inches(3), Inches(3))
c2.fill.solid(); c2.fill.fore_color.rgb = RGBColor(0x04, 0x3D, 0x2A); c2.line.fill.background()
# Title
t = s.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(9), Inches(1))
p = t.text_frame.paragraphs[0]; p.text = "工作汇报"; p.font.size = Pt(54); p.font.color.rgb = WHITE; p.font.bold = True
t2 = s.shapes.add_textbox(Inches(1.2), Inches(2.8), Inches(9), Inches(0.6))
p2 = t2.text_frame.paragraphs[0]; p2.text = "WORK REPORT"; p2.font.size = Pt(20); p2.font.color.rgb = GREEN; p2.font.italic = True
t3 = s.shapes.add_textbox(Inches(1.2), Inches(3.6), Inches(8), Inches(0.5))
p3 = t3.text_frame.paragraphs[0]; p3.text = "汇报人：[姓名]  |  [日期]"; p3.font.size = Pt(16); p3.font.color.rgb = RGBColor(0x9C, 0xAF, 0xA8)
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(3.45), Inches(2), Inches(0.03))
accent.fill.solid(); accent.fill.fore_color.rgb = GREEN; accent.line.fill.background()

# === SLIDE 2: KPI Overview ===
s = prs.slides.add_slide(BLANK); add_bg(s); add_footer(s, 2)
add_title(s, "关键指标概览", "KPIs at a Glance")
kpis = [("¥ 850万", "总营收", "↑12%"), ("95.2%", "客户满意度", "↑3%"), ("1,245", "新增客户", "↑28%"), ("18天", "平均交付周期", "↓5天")]
for i, (val, label, trend) in enumerate(kpis):
    x = Inches(0.6 + i * 3.15)
    add_card(s, x, Inches(1.8), Inches(2.9), Inches(2.0), WHITE, RGBColor(0xE5, 0xE7, 0xEB))
    t = s.shapes.add_textbox(x + Inches(0.2), Inches(2.0), Inches(2.5), Inches(0.6))
    t.text_frame.paragraphs[0].text = val; t.text_frame.paragraphs[0].font.size = Pt(30); t.text_frame.paragraphs[0].font.color.rgb = GREEN; t.text_frame.paragraphs[0].font.bold = True
    t2 = s.shapes.add_textbox(x + Inches(0.2), Inches(2.7), Inches(2.5), Inches(0.35))
    t2.text_frame.paragraphs[0].text = label; t2.text_frame.paragraphs[0].font.size = Pt(12); t2.text_frame.paragraphs[0].font.color.rgb = GRAY
    t3 = s.shapes.add_textbox(x + Inches(0.2), Inches(3.1), Inches(2.5), Inches(0.3))
    t3.text_frame.paragraphs[0].text = trend; t3.text_frame.paragraphs[0].font.size = Pt(18); t3.text_frame.paragraphs[0].font.color.rgb = GREEN; t3.text_frame.paragraphs[0].font.bold = True
# Chart placeholder
add_card(s, Inches(0.6), Inches(4.2), Inches(12.1), Inches(2.5), LIGHT_GREEN)
t = s.shapes.add_textbox(Inches(4), Inches(5.2), Inches(5), Inches(0.5))
t.text_frame.paragraphs[0].text = "📊  在此处插入趋势图表（折线图/柱状图）"; t.text_frame.paragraphs[0].font.size = Pt(14); t.text_frame.paragraphs[0].font.color.rgb = GREEN
t.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# === SLIDE 3: Completed Work ===
s = prs.slides.add_slide(BLANK); add_bg(s); add_footer(s, 3)
add_title(s, "重点项目完成情况", "Key Projects Completed")
projects = [("项目 Alpha", "已完成", "100%", "提前 3 天交付，获得客户书面表扬", GREEN), ("项目 Beta", "已完成", "100%", "按计划完成，质量评分 4.8/5.0", GREEN), ("项目 Gamma", "进行中", "75%", "当前进度超前 2 周，预计月底交付", RGBColor(0xF5, 0x9E, 0x0B)), ("项目 Delta", "进行中", "40%", "遇到技术难点，已制定应对方案", RGBColor(0xEF, 0x44, 0x44))]
for i, (name, status, prog, note, color) in enumerate(projects):
    y = Inches(1.6 + i * 1.1)
    add_card(s, Inches(0.5), y, Inches(12.3), Inches(0.9), WHITE, RGBColor(0xE5, 0xE7, 0xEB))
    # Status badge
    badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), y + Inches(0.2), Inches(1.2), Inches(0.5))
    badge.fill.solid(); badge.fill.fore_color.rgb = color; badge.line.fill.background()
    bp = badge.text_frame.paragraphs[0]; bp.text = status; bp.font.size = Pt(11); bp.font.color.rgb = WHITE; bp.font.bold = True; bp.alignment = PP_ALIGN.CENTER
    tp = s.shapes.add_textbox(Inches(2.1), y + Inches(0.15), Inches(3), Inches(0.3))
    tp.text_frame.paragraphs[0].text = name; tp.text_frame.paragraphs[0].font.size = Pt(16); tp.text_frame.paragraphs[0].font.color.rgb = DARK; tp.text_frame.paragraphs[0].font.bold = True
    tpp = s.shapes.add_textbox(Inches(5.5), y + Inches(0.2), Inches(1.5), Inches(0.3))
    tpp.text_frame.paragraphs[0].text = f"进度 {prog}"; tpp.text_frame.paragraphs[0].font.size = Pt(14); tpp.text_frame.paragraphs[0].font.color.rgb = color; tpp.text_frame.paragraphs[0].font.bold = True
    tn = s.shapes.add_textbox(Inches(2.1), y + Inches(0.5), Inches(9), Inches(0.3))
    tn.text_frame.paragraphs[0].text = note; tn.text_frame.paragraphs[0].font.size = Pt(11); tn.text_frame.paragraphs[0].font.color.rgb = GRAY

# === SLIDE 4: Data Analysis ===
s = prs.slides.add_slide(BLANK); add_bg(s); add_footer(s, 4)
add_title(s, "数据分析与洞察", "Data Analysis & Insights")
# Three insight cards
insights = [("用户增长", "月活用户突破 50 万\n环比增长 23%\n同比增长 156%"), ("收入趋势", "MRR 达到 ¥120 万\n环比增长 18%\n续费率 92%"), ("市场拓展", "新进入 3 个城市\n覆盖率达 65%\n品牌搜索量 +45%")]
for i, (title, desc) in enumerate(insights):
    x = Inches(0.5 + i * 4.2)
    add_card(s, x, Inches(1.8), Inches(3.9), Inches(2.5), WHITE, RGBColor(0xE5, 0xE7, 0xEB))
    hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.8), Inches(3.9), Inches(0.5))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = GREEN; hdr.line.fill.background()
    hp = hdr.text_frame.paragraphs[0]; hp.text = title; hp.font.size = Pt(14); hp.font.color.rgb = WHITE; hp.font.bold = True; hp.alignment = PP_ALIGN.CENTER
    dp = s.shapes.add_textbox(x + Inches(0.3), Inches(2.5), Inches(3.3), Inches(1.6)); dp.text_frame.word_wrap = True
    dpp = dp.text_frame.paragraphs[0]; dpp.text = desc; dpp.font.size = Pt(12); dpp.font.color.rgb = DARK; dpp.line_spacing = Pt(22)
# Bottom summary
add_card(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.3), LIGHT_GREEN)
t = s.shapes.add_textbox(Inches(1), Inches(4.9), Inches(11), Inches(1.8)); t.text_frame.word_wrap = True
p = t.text_frame.paragraphs[0]; p.text = "💡 核心洞察："
p.font.size = Pt(14); p.font.color.rgb = DARK_GREEN; p.font.bold = True
p2 = t.text_frame.add_paragraph(); p2.text = "1. 用户增长超预期，需加快服务端扩容以支撑 Q3 新增流量。\n2. 续费率高达 92%，证明产品价值得到认可，可适度提价。\n3. 新市场拓展顺利，建议追加本地化运营投入。\n4. 竞品动作频繁，建议 Q4 推出差异化功能保持领先。"
p2.font.size = Pt(11); p2.font.color.rgb = DARK; p2.line_spacing = Pt(18); p2.space_before = Pt(8)

# === SLIDE 5: Problems & Solutions ===
s = prs.slides.add_slide(BLANK); add_bg(s); add_footer(s, 5)
add_title(s, "问题与解决方案", "Challenges & Solutions")
issues = [("⚠️ 服务器稳定性不足", "近期出现 3 次服务降级", "已完成架构升级方案评审，Q3 实施", "已完成"), ("⚠️ 竞品上线同类功能", "竞品 A 本月推出免费版本", "差异化竞争：聚焦企业级 B 端客户", "进行中"), ("⚠️ 核心员工离职风险", "2 名高级工程师收到猎头邀请", "启动股权激励 + 薪酬调整方案", "计划中")]
for i, (issue, detail, solution, status) in enumerate(issues):
    y = Inches(1.5 + i * 1.5)
    add_card(s, Inches(0.5), y, Inches(12.3), Inches(1.2), WHITE, RGBColor(0xE5, 0xE7, 0xEB))
    t1 = s.shapes.add_textbox(Inches(0.7), y + Inches(0.1), Inches(5), Inches(0.3))
    t1.text_frame.paragraphs[0].text = f"{issue}：{detail}"; t1.text_frame.paragraphs[0].font.size = Pt(13); t1.text_frame.paragraphs[0].font.color.rgb = DARK; t1.text_frame.paragraphs[0].font.bold = True
    t2 = s.shapes.add_textbox(Inches(0.7), y + Inches(0.5), Inches(6), Inches(0.25))
    t2.text_frame.paragraphs[0].text = f"✅ 解决方案：{solution}"; t2.text_frame.paragraphs[0].font.size = Pt(11); t2.text_frame.paragraphs[0].font.color.rgb = GREEN
    sb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.5), y + Inches(0.35), Inches(1.5), Inches(0.45))
    sb.fill.solid(); sb.fill.fore_color.rgb = GREEN; sb.line.fill.background()
    sp = sb.text_frame.paragraphs[0]; sp.text = status; sp.font.size = Pt(10); sp.font.color.rgb = WHITE; sp.font.bold = True; sp.alignment = PP_ALIGN.CENTER

# === SLIDE 6: Next Quarter Plan ===
s = prs.slides.add_slide(BLANK); add_bg(s); add_footer(s, 6)
add_title(s, "下阶段工作计划", "Next Quarter Plan")
plans = [("产品迭代", ["上线 AI 智能推荐功能", "完成移动端适配", "推出企业版高级套餐"]), ("市场拓展", ["启动华东区域推广活动", "参加 3 场行业峰会", "建立 5 个渠道合作伙伴"]), ("团队建设", ["招聘 8 名核心岗位人员", "启动内部培训体系", "完善绩效考核机制"]), ("财务目标", ["Q3 营收目标 ¥500 万", "毛利率提升至 68%", "启动 B 轮融资准备"])]
for i, (cat, items) in enumerate(plans):
    x = Inches(0.5 + i * 3.15)
    add_card(s, x, Inches(1.5), Inches(3.0), Inches(2.2), WHITE, RGBColor(0xE5, 0xE7, 0xEB))
    hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.5), Inches(3.0), Inches(0.45))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = GREEN; hdr.line.fill.background()
    hp = hdr.text_frame.paragraphs[0]; hp.text = cat; hp.font.size = Pt(13); hp.font.color.rgb = WHITE; hp.font.bold = True; hp.alignment = PP_ALIGN.CENTER
    for j, item in enumerate(items):
        tb = s.shapes.add_textbox(x + Inches(0.2), Inches(2.15 + j * 0.5), Inches(2.6), Inches(0.35))
        tb.text_frame.paragraphs[0].text = f"☐  {item}"; tb.text_frame.paragraphs[0].font.size = Pt(11); tb.text_frame.paragraphs[0].font.color.rgb = DARK
# Timeline bar
add_card(s, Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.8), LIGHT_GREEN)
t = s.shapes.add_textbox(Inches(1.5), Inches(4.3), Inches(10), Inches(2.3)); t.text_frame.word_wrap = True
p = t.text_frame.paragraphs[0]; p.text = "📅 时间规划"; p.font.size = Pt(14); p.font.color.rgb = DARK_GREEN; p.font.bold = True
p2 = t.text_frame.add_paragraph(); p2.text = "\n7月 · 产品迭代启动 | 8月 · 市场推广第一波 | 9月 · 团队扩充到位\n10月 · 企业版发布 | 11月 · 融资材料准备 | 12月 · 年度目标冲刺"
p2.font.size = Pt(12); p2.font.color.rgb = DARK; p2.line_spacing = Pt(24); p2.space_before = Pt(10)

# === SLIDE 7: Resource Needs ===
s = prs.slides.add_slide(BLANK); add_bg(s); add_footer(s, 7)
add_title(s, "需要支持与资源", "Resources & Support Needed")
resources = [("人力需求", ["高级前端工程师 ×2", "数据分析师 ×1", "客户成功经理 ×1"]), ("预算需求", ["市场推广费 ¥50万", "服务器扩容 ¥20万", "团队培训费 ¥8万"]), ("跨部门支持", ["法务：合同模板审核", "财务：Q3 预算快速审批", "HR：核心岗位优先招聘"])]
for i, (cat, items) in enumerate(resources):
    x = Inches(0.5 + i * 4.2)
    add_card(s, x, Inches(1.6), Inches(3.9), Inches(2.8), WHITE, RGBColor(0xE5, 0xE7, 0xEB))
    hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.6), Inches(3.9), Inches(0.45))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = GREEN; hdr.line.fill.background()
    hp = hdr.text_frame.paragraphs[0]; hp.text = cat; hp.font.size = Pt(14); hp.font.color.rgb = WHITE; hp.font.bold = True; hp.alignment = PP_ALIGN.CENTER
    for j, item in enumerate(items):
        tb = s.shapes.add_textbox(x + Inches(0.3), Inches(2.25 + j * 0.65), Inches(3.3), Inches(0.45))
        tb.text_frame.paragraphs[0].text = f"▶  {item}"; tb.text_frame.paragraphs[0].font.size = Pt(12); tb.text_frame.paragraphs[0].font.color.rgb = DARK
# Bottom note
add_card(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.2), LIGHT_GREEN)
t = s.shapes.add_textbox(Inches(1), Inches(5.0), Inches(11), Inches(1.8)); t.text_frame.word_wrap = True
p = t.text_frame.paragraphs[0]; p.text = "📌 期望决策"; p.font.size = Pt(14); p.font.color.rgb = DARK_GREEN; p.font.bold = True
p2 = t.text_frame.add_paragraph(); p2.text = "1. 批准 Q3 市场推广预算 ¥50 万\n2. 确认优先招聘岗位顺序（建议：前端 > 数据 > 客户成功）\n3. 协调法务/财务/HR 部门在下周内完成审批流程"
p2.font.size = Pt(12); p2.font.color.rgb = DARK; p2.line_spacing = Pt(20); p2.space_before = Pt(10)

# === SLIDE 8: Thank You ===
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = DARK_GREEN; bg.line.fill.background()
c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-1.5), Inches(5), Inches(5))
c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0x05, 0x5A, 0x3E); c.line.fill.background()
t = s.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10), Inches(1))
p = t.text_frame.paragraphs[0]; p.text = "谢谢"; p.font.size = Pt(60); p.font.color.rgb = WHITE; p.font.bold = True
t2 = s.shapes.add_textbox(Inches(1.5), Inches(3.0), Inches(10), Inches(0.6))
p2 = t2.text_frame.paragraphs[0]; p2.text = "THANK YOU"; p2.font.size = Pt(22); p2.font.color.rgb = GREEN; p2.font.italic = True
t3 = s.shapes.add_textbox(Inches(1.5), Inches(4.0), Inches(8), Inches(1))
t3.text_frame.word_wrap = True
p3 = t3.text_frame.paragraphs[0]; p3.text = "汇报人：[姓名]\n联系方式：[邮箱]  |  [电话]"; p3.font.size = Pt(14); p3.font.color.rgb = RGBColor(0x9C, 0xAF, 0xA8); p3.line_spacing = Pt(22)
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(3.85), Inches(1.5), Inches(0.03))
accent.fill.solid(); accent.fill.fore_color.rgb = GREEN; accent.line.fill.background()

# Save
out = os.path.join(os.path.dirname(__file__), "工作汇报模板.pptx")
prs.save(out)
print(f"Done: {out} ({len(prs.slides)} slides)")
