"""
Generate a professional 2-page Resume/CV PowerPoint template.
Minimal clean design, editable text, print-ready.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
ACCENT = RGBColor(0x1A, 0x56, 0xDB)
DARK = RGBColor(0x1F, 0x29, 0x3A)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE = RGBColor(0xEB, 0xF0, 0xFF)

prs = Presentation()
# Use letter-size for printability
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# === SLIDE 1: CV main page ===
s = prs.slides.add_slide(BLANK)

# Left sidebar
sidebar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(4.2), prs.slide_height)
sidebar.fill.solid(); sidebar.fill.fore_color.rgb = DARK; sidebar.line.fill.background()

# Photo placeholder
photo = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), Inches(0.7), Inches(2.2), Inches(2.2))
photo.fill.solid(); photo.fill.fore_color.rgb = RGBColor(0x37, 0x47, 0x61); photo.line.color.rgb = ACCENT
photo.line.width = Pt(3)

# Name
t = s.shapes.add_textbox(Inches(4.7), Inches(0.7), Inches(8), Inches(0.6))
t.text_frame.paragraphs[0].text = "张  伟"; t.text_frame.paragraphs[0].font.size = Pt(36)
t.text_frame.paragraphs[0].font.color.rgb = DARK; t.text_frame.paragraphs[0].font.bold = True

# Title line
t2 = s.shapes.add_textbox(Inches(4.7), Inches(1.35), Inches(8), Inches(0.35))
t2.text_frame.paragraphs[0].text = "高级全栈开发工程师  ·  10年经验"
t2.text_frame.paragraphs[0].font.size = Pt(15); t2.text_frame.paragraphs[0].font.color.rgb = GRAY

# Contact bar on sidebar
contacts = [("📧", "zhangwei@email.com"), ("📱", "138-XXXX-XXXX"), ("📍", "北京市朝阳区"), ("🔗", "github.com/zhangwei")]
for i, (icon, text) in enumerate(contacts):
    y = Inches(3.2 + i * 0.5)
    tb = s.shapes.add_textbox(Inches(0.4), y, Inches(3.5), Inches(0.35))
    tb.text_frame.paragraphs[0].text = f"{icon}  {text}"; tb.text_frame.paragraphs[0].font.size = Pt(11)
    tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)

# Skills section on sidebar
sect_title = s.shapes.add_textbox(Inches(0.4), Inches(5.3), Inches(3.5), Inches(0.35))
sect_title.text_frame.paragraphs[0].text = "专业技能"; sect_title.text_frame.paragraphs[0].font.size = Pt(14)
sect_title.text_frame.paragraphs[0].font.color.rgb = WHITE; sect_title.text_frame.paragraphs[0].font.bold = True
# Underline
ul = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(5.65), Inches(1.2), Inches(0.03))
ul.fill.solid(); ul.fill.fore_color.rgb = ACCENT; ul.line.fill.background()

skills = ["React / Next.js", "TypeScript / JavaScript", "Python / Node.js", "PostgreSQL / MongoDB", "AWS / Docker / K8s", "AI / LLM 应用开发"]
for i, skill in enumerate(skills):
    tb = s.shapes.add_textbox(Inches(0.6), Inches(5.9 + i * 0.32), Inches(3.2), Inches(0.28))
    tb.text_frame.paragraphs[0].text = f"▸  {skill}"; tb.text_frame.paragraphs[0].font.size = Pt(10)
    tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xCA, 0xD5, 0xE0)

# Right side: Work Experience
# Section header
sh = s.shapes.add_textbox(Inches(4.7), Inches(2.0), Inches(8), Inches(0.35))
sh.text_frame.paragraphs[0].text = "工作经历"; sh.text_frame.paragraphs[0].font.size = Pt(18)
sh.text_frame.paragraphs[0].font.color.rgb = DARK; sh.text_frame.paragraphs[0].font.bold = True
ul2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.7), Inches(2.35), Inches(1.5), Inches(0.03))
ul2.fill.solid(); ul2.fill.fore_color.rgb = ACCENT; ul2.line.fill.background()

jobs = [
    ("字节跳动", "高级前端工程师", "2022.03 - 至今", [
        "带领 5 人团队负责电商平台前端架构，日活 500 万+",
        "主导前端性能优化，首屏加载时间从 3.2s 降至 1.1s",
        "搭建组件库（200+ 组件），覆盖 8 个业务线",
    ]),
    ("阿里巴巴", "前端开发工程师", "2019.07 - 2022.02", [
        "参与双11大促前端开发，支撑峰值 QPS 10万+",
        "开发内部 DevOps 平台，部署效率提升 60%",
        "获得年度最佳新人奖（Top 5%）",
    ]),
    ("美团", "初级开发工程师", "2016.09 - 2019.06", [
        "负责商家端后台管理系统开发与维护",
        "重构订单模块，代码量减少 40%，性能提升 3x",
    ]),
]

for i, (company, role, period, duties) in enumerate(jobs):
    y = Inches(2.55 + i * 1.55)
    # Company
    tc = s.shapes.add_textbox(Inches(4.7), y, Inches(5), Inches(0.3))
    tc.text_frame.paragraphs[0].text = company; tc.text_frame.paragraphs[0].font.size = Pt(14)
    tc.text_frame.paragraphs[0].font.color.rgb = DARK; tc.text_frame.paragraphs[0].font.bold = True
    # Role + Period
    tr = s.shapes.add_textbox(Inches(4.7), y + Inches(0.3), Inches(5), Inches(0.22))
    tr.text_frame.paragraphs[0].text = f"{role}  |  {period}"; tr.text_frame.paragraphs[0].font.size = Pt(10)
    tr.text_frame.paragraphs[0].font.color.rgb = GRAY
    # Duties
    for j, duty in enumerate(duties):
        td = s.shapes.add_textbox(Inches(4.9), y + Inches(0.55 + j * 0.22), Inches(7.5), Inches(0.22))
        td.text_frame.paragraphs[0].text = f"•  {duty}"; td.text_frame.paragraphs[0].font.size = Pt(10)
        td.text_frame.paragraphs[0].font.color.rgb = DARK

# === SLIDE 2: Education + Projects ===
s2 = prs.slides.add_slide(BLANK)

# Left sidebar
sidebar2 = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(4.2), prs.slide_height)
sidebar2.fill.solid(); sidebar2.fill.fore_color.rgb = DARK; sidebar2.line.fill.background()

# Languages sidebar
st2 = s2.shapes.add_textbox(Inches(0.4), Inches(0.7), Inches(3.5), Inches(0.3))
st2.text_frame.paragraphs[0].text = "语言能力"; st2.text_frame.paragraphs[0].font.size = Pt(14)
st2.text_frame.paragraphs[0].font.color.rgb = WHITE; st2.text_frame.paragraphs[0].font.bold = True
ul3 = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.05), Inches(1.2), Inches(0.03))
ul3.fill.solid(); ul3.fill.fore_color.rgb = ACCENT; ul3.line.fill.background()

langs = [("中文", "母语"), ("英语", "流利（雅思 7.0）"), ("日语", "日常沟通")]
for i, (lang, level) in enumerate(langs):
    tb = s2.shapes.add_textbox(Inches(0.6), Inches(1.3 + i * 0.5), Inches(3.2), Inches(0.35))
    tb.text_frame.paragraphs[0].text = f"{lang}：{level}"; tb.text_frame.paragraphs[0].font.size = Pt(11)
    tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xCA, 0xD5, 0xE0)

# Certificates sidebar
st3 = s2.shapes.add_textbox(Inches(0.4), Inches(3.0), Inches(3.5), Inches(0.3))
st3.text_frame.paragraphs[0].text = "证书与荣誉"; st3.text_frame.paragraphs[0].font.size = Pt(14)
st3.text_frame.paragraphs[0].font.color.rgb = WHITE; st3.text_frame.paragraphs[0].font.bold = True
ul4 = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(3.35), Inches(1.2), Inches(0.03))
ul4.fill.solid(); ul4.fill.fore_color.rgb = ACCENT; ul4.line.fill.background()

certs = ["AWS Solutions Architect", "Google Cloud Professional", "PMP 项目管理认证", "字节跳动年度最佳技术奖", "阿里巴巴双11技术先锋"]
for i, c in enumerate(certs):
    tb = s2.shapes.add_textbox(Inches(0.6), Inches(3.6 + i * 0.32), Inches(3.2), Inches(0.28))
    tb.text_frame.paragraphs[0].text = f"▸  {c}"; tb.text_frame.paragraphs[0].font.size = Pt(10)
    tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xCA, 0xD5, 0xE0)

# Right side: Education
sh2 = s2.shapes.add_textbox(Inches(4.7), Inches(0.7), Inches(8), Inches(0.35))
sh2.text_frame.paragraphs[0].text = "教育背景"; sh2.text_frame.paragraphs[0].font.size = Pt(18)
sh2.text_frame.paragraphs[0].font.color.rgb = DARK; sh2.text_frame.paragraphs[0].font.bold = True
ul5 = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.7), Inches(1.05), Inches(1.5), Inches(0.03))
ul5.fill.solid(); ul5.fill.fore_color.rgb = ACCENT; ul5.line.fill.background()

edu = [("清华大学", "计算机科学与技术 · 硕士", "2014 - 2016", "GPA 3.8/4.0，获国家奖学金"), ("浙江大学", "软件工程 · 学士", "2010 - 2014", "GPA 3.6/4.0，校级优秀毕业生")]
for i, (school, major, period, note) in enumerate(edu):
    y = Inches(1.25 + i * 1.2)
    ts = s2.shapes.add_textbox(Inches(4.7), y, Inches(5), Inches(0.28))
    ts.text_frame.paragraphs[0].text = school; ts.text_frame.paragraphs[0].font.size = Pt(14); ts.text_frame.paragraphs[0].font.color.rgb = DARK; ts.text_frame.paragraphs[0].font.bold = True
    tm = s2.shapes.add_textbox(Inches(4.7), y + Inches(0.3), Inches(5), Inches(0.22))
    tm.text_frame.paragraphs[0].text = f"{major}  |  {period}"; tm.text_frame.paragraphs[0].font.size = Pt(10); tm.text_frame.paragraphs[0].font.color.rgb = GRAY
    tn = s2.shapes.add_textbox(Inches(4.7), y + Inches(0.55), Inches(7), Inches(0.22))
    tn.text_frame.paragraphs[0].text = note; tn.text_frame.paragraphs[0].font.size = Pt(10); tn.text_frame.paragraphs[0].font.color.rgb = DARK

# Right side: Projects
sh3 = s2.shapes.add_textbox(Inches(4.7), Inches(3.7), Inches(8), Inches(0.35))
sh3.text_frame.paragraphs[0].text = "个人项目"; sh3.text_frame.paragraphs[0].font.size = Pt(18)
sh3.text_frame.paragraphs[0].font.color.rgb = DARK; sh3.text_frame.paragraphs[0].font.bold = True
ul6 = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.7), Inches(4.05), Inches(1.5), Inches(0.03))
ul6.fill.solid(); ul6.fill.fore_color.rgb = ACCENT; ul6.line.fill.background()

projects = [
    ("AI周报生成器", "Next.js + Anthropic API", "10万+ 月活用户，付费转化率 8%"),
    ("开源 UI 组件库", "React + TypeScript", "GitHub 2,000+ Stars，被 50+ 项目引用"),
    ("自动化运维平台", "Python + Docker", "内部使用，节省 200+ 人天/年"),
]
for i, (name, tech, desc) in enumerate(projects):
    y = Inches(4.25 + i * 0.8)
    tp = s2.shapes.add_textbox(Inches(4.7), y, Inches(5), Inches(0.25))
    tp.text_frame.paragraphs[0].text = name; tp.text_frame.paragraphs[0].font.size = Pt(13); tp.text_frame.paragraphs[0].font.color.rgb = DARK; tp.text_frame.paragraphs[0].font.bold = True
    tt = s2.shapes.add_textbox(Inches(4.7), y + Inches(0.28), Inches(5), Inches(0.2))
    tt.text_frame.paragraphs[0].text = f"技术栈：{tech}"; tt.text_frame.paragraphs[0].font.size = Pt(10); tt.text_frame.paragraphs[0].font.color.rgb = GRAY
    td = s2.shapes.add_textbox(Inches(4.7), y + Inches(0.5), Inches(7), Inches(0.2))
    td.text_frame.paragraphs[0].text = desc; td.text_frame.paragraphs[0].font.size = Pt(10); td.text_frame.paragraphs[0].font.color.rgb = DARK

# Bottom accent line
bl = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.7), Inches(6.9), Inches(7.5), Inches(0.02))
bl.fill.solid(); bl.fill.fore_color.rgb = RGBColor(0xE5, 0xE7, 0xEB); bl.line.fill.background()

# Save
out = os.path.join(os.path.dirname(__file__), "专业简历模板.pptx")
prs.save(out)
print(f"Done: {out} ({len(prs.slides)} slides)")
